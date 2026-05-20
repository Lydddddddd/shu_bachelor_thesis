from __future__ import annotations

from pathlib import Path
from math import ceil

import fitz
from PIL import Image, ImageEnhance, ImageFilter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "figures"
ASSET_DIR = ROOT / "build" / "ppt_assets"
OUT = Path("/Users/liyaodong/Desktop/最终答辩_多模态混合专家模型的推理优化.pptx")

SLIDE_W = 13.333
SLIDE_H = 7.5

BLUE = "0B5CAD"
NAVY = "102A43"
INK = "172033"
MUTED = "5E6B7A"
PALE = "F5F8FC"
LINE = "D7E2EE"
TEAL = "12A7A3"
CORAL = "E65F5C"
YELLOW = "F6C85F"
GREEN = "3BAE66"
WHITE = "FFFFFF"

FONT_HEAD = "PingFang SC"
FONT_BODY = "PingFang SC"


def set_typeface(run, font_name: str) -> None:
    run.font.name = font_name
    r_pr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = r_pr.find(qn(tag))
        if el is None:
            el = OxmlElement(tag)
            r_pr.append(el)
        el.set("typeface", font_name)


def rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.strip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def ensure_assets() -> None:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    for name in ["shulogo.pdf", "shublack.pdf"]:
        src = FIG / name
        dst = ASSET_DIR / f"{src.stem}.png"
        if dst.exists():
            continue
        doc = fitz.open(src)
        pix = doc[0].get_pixmap(matrix=fitz.Matrix(3, 3), alpha=True)
        pix.save(dst)


def add_textbox(slide, text, x, y, w, h, font=FONT_BODY, size=18, color=INK,
                bold=False, align="left", valign="top", fill=None, margin=0.05):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }[valign]
    p = tf.paragraphs[0]
    p.alignment = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }[align]
    run = p.add_run()
    run.text = text
    set_typeface(run, font)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
        shape.line.color.rgb = rgb(fill)
    return shape


def add_run(p, text, size=16, color=INK, bold=False, font=FONT_BODY):
    run = p.add_run()
    run.text = text
    set_typeface(run, font)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return run


def rect(slide, x, y, w, h, fill=WHITE, line=LINE, radius=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(1)
    return shape


def pill(slide, text, x, y, w, h, fill=BLUE, color=WHITE, size=12):
    s = rect(slide, x, y, w, h, fill=fill, line=fill, radius=True)
    s.text_frame.clear()
    s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = s.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_run(p, text, size=size, color=color, bold=True)
    return s


def add_bullets(slide, items, x, y, w, h, size=17, color=INK, bullet_color=BLUE,
                gap=0.17, lead=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    if lead:
        p = tf.paragraphs[0]
        add_run(p, lead, size=size + 1, color=color, bold=True)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 and not lead else tf.add_paragraph()
        p.space_after = Pt(gap * 18)
        add_run(p, "●  ", size=size, color=bullet_color, bold=True)
        if isinstance(item, tuple):
            head, body = item
            add_run(p, head, size=size, color=color, bold=True)
            add_run(p, body, size=size, color=color)
        else:
            add_run(p, str(item), size=size, color=color)
    return box


def add_header(slide, no, title, subtitle=None):
    add_textbox(slide, f"{no:02d}", 0.42, 0.28, 0.48, 0.28, size=10, color=BLUE, bold=True, align="center")
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.95), Inches(0.36), Inches(0.55), Inches(0.035)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = rgb(BLUE)
    slide.shapes[-1].line.color.rgb = rgb(BLUE)
    add_textbox(slide, title, 1.62, 0.21, 7.4, 0.46, size=24, color=INK, bold=True)
    if subtitle:
        add_textbox(slide, subtitle, 9.3, 0.28, 3.35, 0.34, size=10.5, color=MUTED, align="right")
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.42), Inches(0.83), Inches(12.5), Inches(0.012)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = rgb(LINE)
    slide.shapes[-1].line.color.rgb = rgb(LINE)


def add_footer(slide, page_no):
    add_textbox(slide, "多模态混合专家模型的推理优化", 0.42, 7.08, 4.8, 0.22, size=8.5, color="7A8794")
    add_textbox(slide, str(page_no), 12.45, 7.08, 0.45, 0.22, size=8.5, color="7A8794", align="right")


def add_picture_fit(slide, path, x, y, w, h, fit="contain"):
    path = Path(path)
    with Image.open(path) as im:
        iw, ih = im.size
    box_ratio = w / h
    img_ratio = iw / ih
    pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    if fit == "cover":
        if img_ratio > box_ratio:
            crop = 1 - box_ratio / img_ratio
            pic.crop_left = crop / 2
            pic.crop_right = crop / 2
        else:
            crop = 1 - img_ratio / box_ratio
            pic.crop_top = crop / 2
            pic.crop_bottom = crop / 2
    return pic


def create_soft_bg(src: Path, dst: Path) -> None:
    if dst.exists():
        return
    img = Image.open(src).convert("RGB")
    img = img.resize((1600, int(1600 * img.height / img.width)))
    img = ImageEnhance.Color(img).enhance(0.45)
    img = ImageEnhance.Brightness(img).enhance(1.14)
    img = img.filter(ImageFilter.GaussianBlur(radius=1.3))
    img.save(dst)


def make_table(slide, data, x, y, w, h, col_widths=None, header_fill=BLUE):
    rows, cols = len(data), len(data[0])
    table = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    if col_widths:
        for j, width in enumerate(col_widths):
            table.columns[j].width = Inches(width)
    for i, row in enumerate(data):
        for j, value in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(value)
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(header_fill if i == 0 else WHITE)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    set_typeface(run, FONT_BODY)
                    run.font.size = Pt(10.5 if rows > 5 else 12)
                    run.font.bold = i == 0
                    run.font.color.rgb = rgb(WHITE if i == 0 else INK)
    return table


def add_metric_card(slide, label, value, note, x, y, w, h, color=BLUE):
    rect(slide, x, y, w, h, fill=WHITE, line=LINE, radius=True)
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = rgb(color)
    slide.shapes[-1].line.color.rgb = rgb(color)
    add_textbox(slide, label, x + 0.22, y + 0.22, w - 0.35, 0.28, size=11, color=MUTED, bold=True)
    add_textbox(slide, value, x + 0.22, y + 0.52, w - 0.35, 0.55, size=24, color=color, bold=True)
    add_textbox(slide, note, x + 0.22, y + 1.07, w - 0.35, 0.35, size=10.5, color=MUTED)


def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def main() -> None:
    ensure_assets()
    soft_arch = ASSET_DIR / "vlmoe_architecture_soft.png"
    create_soft_bg(FIG / "vlmoe_architecture.png", soft_arch)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    # 1 Cover
    slide = add_slide(prs)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill="F7FBFF", line="F7FBFF")
    add_picture_fit(slide, soft_arch, 0, 0.3, SLIDE_W, 4.5, fit="cover")
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill="F7FBFF", line="F7FBFF").fill.transparency = 18
    add_picture_fit(slide, ASSET_DIR / "shublack.png", 0.55, 0.35, 2.2, 0.65, fit="contain")
    add_textbox(slide, "多模态混合专家模型的推理优化", 0.78, 1.65, 8.8, 0.82, size=33, color=NAVY, bold=True)
    add_textbox(slide, "Efficient Inference for Vision-Language Mixture-of-Experts Models", 0.82, 2.55, 8.5, 0.35, size=15, color=BLUE, bold=False)
    rect(slide, 0.82, 3.16, 4.3, 0.06, fill=TEAL, line=TEAL)
    info = [("答辩人", "李耀东"), ("导师", "滕中梅"), ("专业", "计算机科学与技术"), ("日期", "2026年5月")]
    for i, (k, v) in enumerate(info):
        add_textbox(slide, k, 0.88 + i * 1.65, 6.36, 0.58, 0.26, size=10.5, color=MUTED, bold=True)
        add_textbox(slide, v, 0.88 + i * 1.65, 6.72, 1.45, 0.28, size=12.5, color=INK)
    pill(slide, "本科毕业论文答辩", 10.2, 0.56, 2.25, 0.35, fill=BLUE, size=12)

    # 2 Outline
    slide = add_slide(prs)
    add_header(slide, 2, "汇报提纲", "从中期方案到最终闭环")
    parts = [
        ("01", "研究背景与问题定义", "VL-MoE 推理链路的双重瓶颈"),
        ("02", "负载分析与优化目标", "视觉 Token 冗余 + 专家跨层可预测"),
        ("03", "方法设计与系统实现", "SAT-ToMe / Pipeline / Global Predictor"),
        ("04", "实验结果与效果分析", "TTFT、Recall、真实 offload 微基准"),
        ("05", "总结、不足与展望", "贡献归纳与后续工程化方向"),
    ]
    for i, (num, title, sub) in enumerate(parts):
        y = 1.28 + i * 1.0
        pill(slide, num, 1.05, y, 0.58, 0.42, fill=[BLUE, TEAL, CORAL, YELLOW, NAVY][i], size=13)
        add_textbox(slide, title, 1.85, y - 0.02, 3.2, 0.36, size=18, color=INK, bold=True)
        add_textbox(slide, sub, 5.0, y + 0.03, 5.9, 0.28, size=12.5, color=MUTED)
        if i < 4:
            slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.32), Inches(y + 0.44), Inches(0.035), Inches(0.5)).fill.solid()
            slide.shapes[-1].fill.fore_color.rgb = rgb(LINE)
            slide.shapes[-1].line.color.rgb = rgb(LINE)
    add_footer(slide, 2)

    # 3 Background
    slide = add_slide(prs)
    add_header(slide, 3, "研究背景：VL-MoE 把多模态能力与稀疏专家结合", "但推理链路更长、更异构")
    add_picture_fit(slide, FIG / "vlmoe_architecture.png", 0.75, 1.12, 7.4, 2.55, fit="contain")
    rect(slide, 8.55, 1.18, 3.85, 2.4, fill="F8FBFE", line=LINE, radius=True)
    add_bullets(slide, [
        ("视觉前端：", "图像被编码为大量视觉 Token，推理开始前必须完成。"),
        ("稀疏主干：", "MoE 只激活部分专家，但访问集合由 Router 动态决定。"),
        ("部署约束：", "显存受限时专家权重需要 CPU-GPU 搬运。"),
    ], 8.85, 1.45, 3.3, 1.85, size=13.5)
    for i, (label, body, color) in enumerate([
        ("首字延迟", "Visual Encoder 直接影响用户可感知等待", BLUE),
        ("动态调度", "专家集合随层和请求变化，难以静态缓存", CORAL),
        ("资源受限", "大模型部署常受 GPU cache 容量约束", TEAL),
    ]):
        add_metric_card(slide, label, body.split("，")[0], body.split("，")[-1] if "，" in body else "", 0.9 + i * 4.15, 4.5, 3.55, 1.25, color=color)
    add_footer(slide, 3)

    # 4 Profiling
    slide = add_slide(prs)
    add_header(slide, 4, "Profiling 结论：瓶颈不是单点，而是前端与专家侧叠加", "中期分析在最终稿中被量化补全")
    add_picture_fit(slide, FIG / "token_distribution_mmbench.png", 0.65, 1.25, 2.3, 1.82, fit="contain")
    add_picture_fit(slide, FIG / "token_distribution_mme.png", 2.95, 1.25, 2.3, 1.82, fit="contain")
    add_metric_card(slide, "图像 Token 占比", "95.08% / 98.15%", "MMBench / MME", 0.9, 3.25, 4.1, 1.05, color=BLUE)
    add_picture_fit(slide, FIG / "layer_correlation_heatmap.jpeg", 5.8, 1.15, 3.0, 3.15, fit="contain")
    add_picture_fit(slide, FIG / "layer_correlation_stats.png", 8.7, 1.28, 3.65, 2.82, fit="contain")
    add_bullets(slide, [
        "Visual Encoder 在基线 TTFT 中占 20.6%，是首字延迟的关键前端瓶颈。",
        "单层专家使用较均衡，但跨层转移路径稀疏且存在统计集中性。",
        "优化方向由中期的三个尝试收敛为两条主线：视觉侧降冗余，专家侧提前调度。",
    ], 0.85, 5.02, 11.6, 1.1, size=14.5)
    add_footer(slide, 4)

    # 5 Problem model
    slide = add_slide(prs)
    add_header(slide, 5, "问题建模：在精度与显存约束下最小化端到端时延", "把论文主线讲成一个优化目标")
    rect(slide, 0.85, 1.15, 11.8, 1.15, fill=NAVY, line=NAVY, radius=True)
    add_textbox(slide, "L = Tenc + Tfuse + Troute + Texpert + Tcomm + Tdecode", 1.15, 1.42, 8.0, 0.38, size=22, color=WHITE, bold=True)
    add_textbox(slide, "核心压缩项：Tenc 与 Troute + Tcomm", 9.25, 1.48, 2.9, 0.26, size=11.5, color="D8F3FF", bold=True, align="right")
    cards = [
        ("目标一", "降低视觉侧冗余计算", "在效果基本不变的前提下压缩视觉 Token，并通过流水隐藏等待。", BLUE),
        ("目标二", "降低专家访问等待", "利用路由可预测性提前预取高概率专家，减少同步加载。", CORAL),
        ("约束", "精度下降与显存预算可控", "ΔAcc ≤ ε，M ≤ Mbudget，错误预测必须有回退机制。", TEAL),
    ]
    for i, (t, h, b, c) in enumerate(cards):
        x = 0.9 + i * 4.08
        rect(slide, x, 3.05, 3.55, 2.35, fill=WHITE, line=LINE, radius=True)
        pill(slide, t, x + 0.22, 3.28, 0.75, 0.32, fill=c, size=10)
        add_textbox(slide, h, x + 0.22, 3.8, 3.05, 0.36, size=17, color=INK, bold=True)
        add_textbox(slide, b, x + 0.22, 4.35, 3.0, 0.72, size=12.5, color=MUTED)
    add_footer(slide, 5)

    # 6 Overall route
    slide = add_slide(prs)
    add_header(slide, 6, "总体技术路线：视觉侧与专家侧互补优化", "最终稿比中期稿新增了真实 offload 验证闭环")
    steps = [
        ("Profiling", "定位视觉 Token 与专家路由瓶颈", BLUE),
        ("SAT-ToMe", "文本引导自适应 Token 压缩", TEAL),
        ("Pipeline", "视觉流与主流异步重叠", YELLOW),
        ("Global Predictor", "layer-0 条件化专家预测", CORAL),
        ("Prefetch Scheduler", "受限预取与缓存回退", NAVY),
    ]
    for i, (h, b, c) in enumerate(steps):
        x = 0.62 + i * 2.52
        rect(slide, x, 2.25, 2.05, 1.35, fill=WHITE, line=LINE, radius=True)
        add_textbox(slide, h, x + 0.18, 2.48, 1.68, 0.28, size=13.5, color=c, bold=True, align="center")
        add_textbox(slide, b, x + 0.18, 2.9, 1.68, 0.42, size=10.5, color=MUTED, align="center")
        if i < len(steps) - 1:
            add_textbox(slide, "→", x + 2.05, 2.63, 0.45, 0.35, size=26, color="9AA8B5", align="center")
    make_table(slide, [
        ["中期稿", "最终答辩应修改为"],
        ["阶段性成果：已实现/后续计划", "最终成果：方法、系统实现、实验闭环"],
        ["ExpertPredictor 只给 +1/+2/+3 R@8", "全局预测器 best mean recall 0.7219"],
        ["专家缓存结果偏模拟/受限", "真实 CPU-GPU offload 微基准 2.99x / 2.66x"],
    ], 1.2, 4.58, 10.9, 1.28, col_widths=[3.4, 7.5], header_fill=NAVY)
    add_footer(slide, 6)

    # 7 SAT-ToMe framework
    slide = add_slide(prs)
    add_header(slide, 7, "方法一：语义感知视觉 Token 压缩 SAT-ToMe", "不是固定裁剪，而是由文本问题引导保留区域")
    add_picture_fit(slide, FIG / "sat_tome_framework.png", 0.65, 1.15, 7.65, 3.6, fit="contain")
    rect(slide, 8.75, 1.18, 3.55, 3.45, fill="F8FBFE", line=LINE, radius=True)
    add_bullets(slide, [
        ("语义相关性：", "文本特征投影到视觉空间，计算每个视觉 Token 的相关分数。"),
        ("自适应比例：", "用语义分布熵判断信息集中度，动态决定保留数量 K。"),
        ("工程实现：", "当前版本采用轻量 Top-K 保留，并同步更新位置编码与变长序列元信息。"),
    ], 9.05, 1.52, 3.0, 2.58, size=13)
    add_picture_fit(slide, FIG / "semantic_token_guidance_example.png", 0.9, 5.0, 11.55, 1.3, fit="contain")
    add_footer(slide, 7)

    # 8 Formula and adaptive logic
    slide = add_slide(prs)
    add_header(slide, 8, "SAT-ToMe 的核心：相关性打分 + 熵自适应压缩", "回答“保留哪些 Token、保留多少 Token”")
    formulas = [
        ("文本投影", "wt = Φ(Tfeat)"),
        ("相关性打分", "si = cos(wt, vi)"),
        ("归一化熵", "H(S) = -1/logN · Σ pi log pi"),
        ("保留数量", "K = max(Kmin, floor(N · (1 - Radaptive)))"),
    ]
    for i, (h, f) in enumerate(formulas):
        x = 0.88 + (i % 2) * 5.75
        y = 1.25 + (i // 2) * 1.48
        rect(slide, x, y, 5.15, 1.05, fill=WHITE, line=LINE, radius=True)
        add_textbox(slide, h, x + 0.22, y + 0.2, 1.35, 0.28, size=12, color=BLUE, bold=True)
        add_textbox(slide, f, x + 1.6, y + 0.17, 3.25, 0.36, size=17, color=INK, bold=True, font="Aptos")
    add_bullets(slide, [
        ("熵低：", "语义集中在少数区域，压缩可以更激进。"),
        ("熵高：", "问题需要更完整场景，系统自动降低压缩强度。"),
        ("阈值保护：", "设置 Kmin 避免极端样本中视觉骨架信息过少。"),
    ], 1.0, 4.62, 5.0, 1.4, size=14)
    add_metric_card(slide, "平均视觉 Token 保留比例", "63.7%", "后续实验在同一保留比例下比较效果", 7.15, 4.55, 4.7, 1.28, color=TEAL)
    add_footer(slide, 8)

    # 9 Pipeline
    slide = add_slide(prs)
    add_header(slide, 9, "方法二：面向实时推理的双流异步流水", "减少计算量之外，继续隐藏前端等待")
    add_picture_fit(slide, FIG / "encoder_pipeline.png", 0.8, 1.1, 11.6, 3.25, fit="contain")
    for i, (h, b, c) in enumerate([
        ("Stage 1", "视觉预处理与前若干视觉层", BLUE),
        ("SAT-ToMe", "文本引导压缩并重建变长序列", TEAL),
        ("Stage 2", "剩余视觉层与 DeepStack 特征", CORAL),
        ("Sync", "仅在依赖交换点 wait_stream", NAVY),
    ]):
        x = 0.95 + i * 3.02
        rect(slide, x, 5.0, 2.55, 0.95, fill=WHITE, line=LINE, radius=True)
        add_textbox(slide, h, x + 0.18, 5.16, 2.15, 0.24, size=13.5, color=c, bold=True, align="center")
        add_textbox(slide, b, x + 0.18, 5.5, 2.15, 0.24, size=10.2, color=MUTED, align="center")
    add_footer(slide, 9)

    # 10 Predictor
    slide = add_slide(prs)
    add_header(slide, 10, "方法三：layer-0 条件化的全局专家路由预测器", "把专家访问从“发生后处理”前移到“访问前准备”")
    add_picture_fit(slide, FIG / "fig4_2_data_pipeline.png", 0.65, 1.08, 5.8, 2.2, fit="contain")
    add_picture_fit(slide, FIG / "fig4_3_dataset_stats.png", 0.8, 3.42, 5.5, 2.0, fit="contain")
    rect(slide, 7.0, 1.18, 5.45, 4.25, fill="F8FBFE", line=LINE, radius=True)
    add_bullets(slide, [
        ("数据：", "重新采集 Qwen3-VL fresh 路由 trace，共 8347 个 token。"),
        ("输入：", "layer-0 router 概率、Top-k 专家、token 位置、input id、模态标识、路由统计量。"),
        ("结构：", "共享 MLP backbone + 多层输出头，直接预测后续层专家 logits。"),
        ("目标：", "多标签 BCEWithLogits，优化未来层专家集合 recall。"),
    ], 7.35, 1.55, 4.65, 3.28, size=13)
    add_footer(slide, 10)

    # 11 Scheduler
    slide = add_slide(prs)
    add_header(slide, 11, "方法四：预测驱动的受限预取与缓存调度", "预测结果不能直接等价为搬运指令")
    add_picture_fit(slide, FIG / "fig4_1_offload_architecture.png", 0.65, 1.1, 11.9, 2.45, fit="contain")
    blocks = [
        ("候选专家", "预测未来层高概率专家集合", BLUE),
        ("价值评分", "置信度 × 层距衰减 × margin", TEAL),
        ("预算约束", "每层 key 上限 + 全局预取预算", YELLOW),
        ("保护回退", "保护当前专家，未命中则同步加载", CORAL),
    ]
    for i, (h, b, c) in enumerate(blocks):
        x = 0.85 + i * 3.05
        rect(slide, x, 4.42, 2.55, 1.15, fill=WHITE, line=LINE, radius=True)
        pill(slide, str(i + 1), x + 0.22, 4.62, 0.35, 0.3, fill=c, size=10)
        add_textbox(slide, h, x + 0.72, 4.6, 1.55, 0.26, size=13.5, color=INK, bold=True)
        add_textbox(slide, b, x + 0.22, 5.02, 2.1, 0.28, size=10.5, color=MUTED, align="center")
    add_footer(slide, 11)

    # 12 Experiment setup
    slide = add_slide(prs)
    add_header(slide, 12, "系统实现与实验设置", "最终答辩强调可复现实验原型")
    make_table(slide, [
        ["项目", "配置"],
        ["硬件", "4 × NVIDIA A6000 GPU，单卡 48GB"],
        ["模型", "Qwen3-VL-30B-A3B-Instruct"],
        ["软件", "Python 3.10 / CUDA 12.4 / PyTorch 2.6.0 / Transformers 4.57.1 / vLLM 0.10.2"],
        ["视觉侧数据", "MMBench 与 MME 固定子集，各 110 条含图像样本"],
        ["专家侧数据", "fresh route trace：8347 tokens，训练 7268 / 验证 1079"],
        ["offload 微基准", "真实 CPU-GPU expert 搬运 + GPU expert MLP 计算"],
    ], 0.82, 1.2, 11.9, 4.18, col_widths=[2.25, 9.65], header_fill=BLUE)
    add_bullets(slide, [
        "端到端 generate 原型可运行，但主结果采用更稳定、可解释的真实 expert 搬运微基准。",
        "所有策略在同一环境、同一 trace 顺序上比较，报告相对变化而不是绝对硬件峰值。",
    ], 0.92, 5.72, 11.3, 0.8, size=13.5)
    add_footer(slide, 12)

    # 13 Quality
    slide = add_slide(prs)
    add_header(slide, 13, "实验一：视觉压缩保持任务效果稳定", "同一保留比例下，语义引导明显优于随机裁剪和固定合并")
    make_table(slide, [
        ["方法", "保留比例", "MMBench", "MME", "输出一致率"],
        ["Baseline", "100.0%", "1.000", "1.000", "100.0%"],
        ["Random pruning", "63.7%", "0.948", "0.936", "89.1%"],
        ["Static ToMe", "63.7%", "0.982", "0.976", "96.8%"],
        ["SAT-ToMe", "63.7%", "0.997", "0.995", "99.1%"],
        ["SAT-ToMe + Pipeline", "63.7%", "0.997", "0.995", "99.1%"],
    ], 0.9, 1.25, 8.25, 3.55, col_widths=[2.7, 1.35, 1.35, 1.1, 1.75], header_fill=NAVY)
    add_metric_card(slide, "归一化任务分数", "> 0.995", "MMBench 与 MME 均保持稳定", 9.65, 1.55, 2.75, 1.25, color=TEAL)
    add_metric_card(slide, "输出一致率", "99.1%", "压缩后与 baseline 高度一致", 9.65, 3.15, 2.75, 1.25, color=BLUE)
    add_bullets(slide, [
        "随机裁剪说明视觉 Token 不能被无差别删除。",
        "Pipeline 不改变语义计算图，效果与 SAT-ToMe 保持一致。",
    ], 1.0, 5.5, 11.2, 0.75, size=13.5)
    add_footer(slide, 13)

    # 14 TTFT
    slide = add_slide(prs)
    add_header(slide, 14, "实验二：Visual Encoder 可见耗时降至 800ms", "整体 TTFT 获得约 1.63x 加速")
    add_picture_fit(slide, FIG / "encoder_ttft.png", 0.9, 1.15, 5.5, 3.65, fit="contain")
    make_table(slide, [
        ["方法", "TTFT/ms", "Encoder/ms", "占比", "TTFT 加速", "Encoder 加速"],
        ["Baseline", "41682", "8572", "20.6%", "1.00x", "1.00x"],
        ["SAT-ToMe", "25738", "1088", "4.2%", "1.62x", "7.88x"],
        ["SAT-ToMe + Pipeline", "25512", "800", "3.1%", "1.63x", "10.72x"],
    ], 6.65, 1.45, 5.75, 2.35, col_widths=[1.55, 0.9, 1.0, 0.72, 0.78, 0.8], header_fill=BLUE)
    add_metric_card(slide, "Encoder 占 TTFT", "20.6% → 3.1%", "前端瓶颈被显著压缩", 6.75, 4.45, 2.65, 1.22, color=CORAL)
    add_metric_card(slide, "Encoder 加速", "10.72x", "Token 合并 + 双流 overlap", 9.75, 4.45, 2.65, 1.22, color=TEAL)
    add_footer(slide, 14)

    # 15 Predictor offline
    slide = add_slide(prs)
    add_header(slide, 15, "实验三：全局预测器学到跨层路由规律", "best mean recall 达到 0.7219")
    add_picture_fit(slide, FIG / "fig4_4_training_curve.png", 0.75, 1.1, 7.25, 3.55, fit="contain")
    make_table(slide, [
        ["视野", "early", "mid", "late"],
        ["Recall", "0.7158", "0.7453", "0.7042"],
    ], 8.35, 1.45, 3.7, 1.35, col_widths=[1.05, 0.88, 0.88, 0.88], header_fill=NAVY)
    add_metric_card(slide, "Best mean recall", "0.7219", "layer-0 条件化全局预测", 8.45, 3.35, 3.45, 1.25, color=BLUE)
    add_bullets(slide, [
        "预测器不是盲猜，而是利用首层 router 概率、Top-k、位置、输入 id 和模态信息。",
        "recall 指标证明未来专家集合存在可利用的跨层结构性。",
    ], 1.0, 5.25, 11.3, 0.85, size=13.5)
    add_footer(slide, 15)

    # 16 Offload benchmark
    slide = add_slide(prs)
    add_header(slide, 16, "实验四：真实 CPU-GPU expert offload 微基准", "预测式调度显著减少同步加载等待")
    add_picture_fit(slide, FIG / "fig5_1_real_offload_compare.png", 0.65, 1.0, 6.75, 3.0, fit="contain")
    make_table(slide, [
        ["setting", "policy", "elapsed_s", "speedup", "cache_hits", "transfer_GB"],
        ["cap96", "demand", "4.283", "1.00x", "0", "9.81"],
        ["cap96", "predictor", "1.432", "2.99x", "104", "14.27"],
        ["cap96", "oracle", "1.354", "3.16x", "111", "14.20"],
        ["cap160", "demand", "4.962", "1.00x", "0", "15.87"],
        ["cap160", "predictor", "1.864", "2.66x", "1134", "17.17"],
        ["cap160", "oracle", "1.571", "3.16x", "1208", "16.55"],
    ], 7.55, 1.05, 5.1, 3.75, col_widths=[0.78, 1.05, 0.82, 0.78, 0.9, 0.77], header_fill=BLUE)
    add_metric_card(slide, "cap96 加速", "2.99x", "predictor 相对 demand", 1.05, 4.98, 2.65, 1.1, color=TEAL)
    add_metric_card(slide, "cap160 加速", "2.66x", "predictor 相对 demand", 4.0, 4.98, 2.65, 1.1, color=CORAL)
    add_bullets(slide, [
        "predictor 传输量通常高于 demand：用受控额外预取换同步等待下降。",
        "oracle 给出上界，说明预测质量与缓存预算仍有优化空间。",
    ], 7.75, 5.12, 4.7, 0.85, size=12.5)
    add_footer(slide, 16)

    # 17 Contributions
    slide = add_slide(prs)
    add_header(slide, 17, "主要工作与创新点", "最终答辩建议用这一页收束贡献")
    contribs = [
        ("1", "完成 VL-MoE 推理负载分析", "从输入 Token、Visual Encoder、专家路由和跨层转移角度定位双重瓶颈。", BLUE),
        ("2", "提出语义感知视觉压缩与流水优化", "文本引导相关性打分 + 熵自适应压缩 + 双流 overlap，降低 TTFT。", TEAL),
        ("3", "构建全局路由预测与调度机制", "fresh route trace、layer-0 条件化预测器、受限预取与缓存回退。", CORAL),
        ("4", "完成真实 offload 微基准验证", "在 CPU-GPU expert 搬运场景下获得 2.99x / 2.66x 加速。", NAVY),
    ]
    for i, (n, h, b, c) in enumerate(contribs):
        y = 1.15 + i * 1.25
        pill(slide, n, 0.9, y + 0.18, 0.48, 0.42, fill=c, size=13)
        rect(slide, 1.65, y, 10.7, 0.88, fill=WHITE, line=LINE, radius=True)
        add_textbox(slide, h, 1.95, y + 0.17, 3.7, 0.3, size=16, color=INK, bold=True)
        add_textbox(slide, b, 5.75, y + 0.2, 6.0, 0.27, size=12.5, color=MUTED)
    add_footer(slide, 17)

    # 18 Limitations
    slide = add_slide(prs)
    add_header(slide, 18, "不足与展望", "把中期的“后续计划”改成更成熟的研究边界")
    left = [
        ("端到端 serving 融合不足", "视觉侧与专家侧分别验证，完整服务框架中仍需更稳定集成。"),
        ("路由数据规模有限", "当前主要覆盖单一模型与实验数据集，跨任务/长上下文泛化需补充。"),
        ("调度仍偏原型系统", "部分逻辑位于 Python 层，缓存替换、传输重叠和预算自适应还可下沉。"),
    ]
    right = [
        ("扩大 trace", "引入更多任务类型、模态比例和上下文长度。"),
        ("下沉 runtime", "降低 Python 调度开销，完善异步传输与动态预算。"),
        ("端到端部署", "结合批处理、KV cache 与多用户负载验证真实 serving 收益。"),
    ]
    add_textbox(slide, "当前不足", 1.0, 1.16, 3.0, 0.34, size=18, color=CORAL, bold=True)
    add_textbox(slide, "后续方向", 7.05, 1.16, 3.0, 0.34, size=18, color=TEAL, bold=True)
    for i, (h, b) in enumerate(left):
        rect(slide, 0.9, 1.75 + i * 1.35, 5.45, 0.96, fill=WHITE, line=LINE, radius=True)
        add_textbox(slide, h, 1.15, 1.91 + i * 1.35, 4.7, 0.25, size=14.5, color=INK, bold=True)
        add_textbox(slide, b, 1.15, 2.28 + i * 1.35, 4.75, 0.28, size=11.5, color=MUTED)
    for i, (h, b) in enumerate(right):
        rect(slide, 6.95, 1.75 + i * 1.35, 5.45, 0.96, fill=WHITE, line=LINE, radius=True)
        add_textbox(slide, h, 7.2, 1.91 + i * 1.35, 4.7, 0.25, size=14.5, color=INK, bold=True)
        add_textbox(slide, b, 7.2, 2.28 + i * 1.35, 4.75, 0.28, size=11.5, color=MUTED)
    add_footer(slide, 18)

    # 19 Closing
    slide = add_slide(prs)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY, line=NAVY)
    add_picture_fit(slide, soft_arch, 0, 0, SLIDE_W, 5.0, fit="cover")
    overlay = rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY, line=NAVY)
    overlay.fill.transparency = 15
    add_textbox(slide, "敬请各位老师批评指正！", 2.15, 2.35, 9.0, 0.75, size=34, color=WHITE, bold=True, align="center")
    add_textbox(slide, "答辩人：李耀东   导师：滕中梅   专业：计算机科学与技术", 3.0, 3.38, 7.3, 0.32, size=14, color="DCEEFF", align="center")
    add_picture_fit(slide, ASSET_DIR / "shulogo.png", 5.95, 4.62, 1.2, 1.45, fit="contain")

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
